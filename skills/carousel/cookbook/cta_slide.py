"""Layout: CTA Slide — Call-to-action slide with button-style shape."""
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor


def add_slide(prs, title="", body="", colors=None, fonts=None, **kwargs):
    """Add a call-to-action slide with a prominent button-style element.

    Centered layout with a bold headline, supporting text, and a
    rounded-rectangle "button" shape with CTA text. Perfect for closing
    slides on LinkedIn carousels and pitch decks.

    Args:
        prs: python-pptx Presentation object
        title: CTA headline (e.g., "Ready to Get Started?")
        body: Supporting text below headline
        colors: dict with hex colors
        fonts: dict with font names
        **kwargs:
            button_text (str): Text on the CTA button (e.g., "Book a Demo")
            url (str): URL for the CTA (added to slide notes)
            handle (str): Social handle or contact info below button

    Returns:
        The created slide object
    """
    if not colors:
        colors = {
            'primary': '1a1a2e',
            'secondary': '16213e',
            'accent': 'e94560',
            'text': 'ffffff',
            'bg': '1a1a2e',
        }
    if not fonts:
        fonts = {'heading': 'Arial Black', 'body': 'Arial'}

    button_text = kwargs.get('button_text', 'Learn More')
    url = kwargs.get('url', '')
    handle = kwargs.get('handle', '')

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(colors['bg'])

    # CTA Headline
    title_left = Inches(1)
    title_top = int(slide_height * 0.22)
    title_width = slide_width - Inches(2)

    txBox = slide.shapes.add_textbox(title_left, title_top, title_width, Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor.from_string(colors['text'])
    p.font.name = fonts['heading']
    p.alignment = PP_ALIGN.CENTER

    # Supporting text
    if body:
        sub_top = int(title_top + Inches(1.6))
        txBox2 = slide.shapes.add_textbox(
            Inches(1.5), sub_top, slide_width - Inches(3), Inches(0.8)
        )
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = body
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor.from_string(colors['text'])
        p2.font.name = fonts['body']
        p2.alignment = PP_ALIGN.CENTER

    # CTA Button (rounded rectangle)
    btn_width = Inches(3.5)
    btn_height = Inches(0.8)
    btn_left = int((slide_width - btn_width) / 2)
    btn_top = int(slide_height * 0.58)

    button = slide.shapes.add_shape(
        5,  # MSO_SHAPE.ROUNDED_RECTANGLE
        btn_left, btn_top, int(btn_width), int(btn_height)
    )
    button.fill.solid()
    button.fill.fore_color.rgb = RGBColor.from_string(colors['accent'])
    button.line.fill.background()

    # Adjust rounding
    button.adjustments[0] = 0.3

    # Button text
    btf = button.text_frame
    btf.word_wrap = False
    bp = btf.paragraphs[0]
    bp.text = button_text
    bp.font.size = Pt(20)
    bp.font.bold = True
    bp.font.color.rgb = RGBColor.from_string(colors['text'])
    bp.font.name = fonts['heading']
    bp.alignment = PP_ALIGN.CENTER

    # Handle / contact info
    if handle:
        handle_top = int(btn_top + btn_height + Inches(0.5))
        handle_box = slide.shapes.add_textbox(
            Inches(1), handle_top, slide_width - Inches(2), Inches(0.5)
        )
        htf = handle_box.text_frame
        hp = htf.paragraphs[0]
        hp.text = handle
        hp.font.size = Pt(16)
        hp.font.color.rgb = RGBColor.from_string(colors['accent'])
        hp.font.name = fonts['body']
        hp.alignment = PP_ALIGN.CENTER

    # Add URL to slide notes
    if url:
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = f"CTA URL: {url}"

    return slide
