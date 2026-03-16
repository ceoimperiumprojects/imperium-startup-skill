"""Layout: Content Single — Single key point with large text."""
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor


def add_slide(prs, title="", body="", colors=None, fonts=None, **kwargs):
    """Add a single key point slide with large centered text.

    Designed for maximum impact — one idea, big text, lots of whitespace.
    Perfect for LinkedIn carousels where each slide delivers one message.

    Args:
        prs: python-pptx Presentation object
        title: Small top label or slide number (optional)
        body: The single key point (keep under 15 words)
        colors: dict with hex colors
        fonts: dict with font names
        **kwargs: Additional parameters (unused)

    Returns:
        The created slide object
    """
    if not colors:
        colors = {
            'primary': '1a1a2e',
            'secondary': '16213e',
            'accent': 'e94560',
            'text': 'ffffff',
            'bg': '1a1a2e',
        }
    if not fonts:
        fonts = {'heading': 'Arial Black', 'body': 'Arial'}

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(colors['bg'])

    # Small label at top (slide number or category)
    if title:
        label_left = Inches(1)
        label_top = Inches(0.8)
        label_width = slide_width - Inches(2)
        label_height = Inches(0.5)

        txBox = slide.shapes.add_textbox(label_left, label_top, label_width, label_height)
        tf = txBox.text_frame

        p = tf.paragraphs[0]
        p.text = title.upper()
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor.from_string(colors['accent'])
        p.font.name = fonts['body']
        p.alignment = PP_ALIGN.CENTER

    # Main point — large, centered, vertically centered
    body_left = Inches(1.2)
    body_top = int(slide_height * 0.28)
    body_width = slide_width - Inches(2.4)
    body_height = int(slide_height * 0.44)

    txBox2 = slide.shapes.add_textbox(body_left, body_top, body_width, body_height)
    tf2 = txBox2.text_frame
    tf2.word_wrap = True

    p2 = tf2.paragraphs[0]
    p2.text = body
    p2.font.size = Pt(36)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor.from_string(colors['text'])
    p2.font.name = fonts['heading']
    p2.alignment = PP_ALIGN.CENTER

    return slide
