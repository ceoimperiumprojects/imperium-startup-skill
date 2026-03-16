"""Layout: Data Table — Clean data table with styled header."""
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def add_slide(prs, title="", body="", colors=None, fonts=None, **kwargs):
    """Add a data table slide with styled header row.

    Creates a table with a colored header row and alternating row
    shading. Supports arbitrary columns and rows.

    Args:
        prs: python-pptx Presentation object
        title: Slide heading
        body: Unused
        colors: dict with hex colors
        fonts: dict with font names
        **kwargs:
            headers (list[str]): Column header labels
            rows (list[list[str]]): Table data rows
            col_widths (list[float]): Column widths in inches (optional)

    Returns:
        The created slide object
    """
    if not colors:
        colors = {
            'primary': '1a1a2e',
            'secondary': '16213e',
            'accent': 'e94560',
            'text': 'ffffff',
            'bg': '1a1a2e',
        }
    if not fonts:
        fonts = {'heading': 'Arial Black', 'body': 'Arial'}

    headers = kwargs.get('headers', ['Column A', 'Column B', 'Column C'])
    rows = kwargs.get('rows', [['Data 1', 'Data 2', 'Data 3']])
    col_widths_inches = kwargs.get('col_widths', None)

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(colors['bg'])

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), slide_width - Inches(1.6), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor.from_string(colors['text'])
    p.font.name = fonts['heading']
    p.alignment = PP_ALIGN.LEFT

    # Table dimensions
    num_cols = len(headers)
    num_rows = len(rows) + 1  # +1 for header row
    table_left = Inches(0.8)
    table_top = Inches(1.6)
    table_width = slide_width - Inches(1.6)
    table_height = min(slide_height - Inches(2.2), Inches(0.5) * num_rows)

    # Create table
    table_shape = slide.shapes.add_table(
        num_rows, num_cols,
        int(table_left), int(table_top),
        int(table_width), int(table_height)
    )
    table = table_shape.table

    # Set column widths
    if col_widths_inches:
        for i, w in enumerate(col_widths_inches):
            if i < num_cols:
                table.columns[i].width = Inches(w)
    else:
        col_width = int(table_width / num_cols)
        for i in range(num_cols):
            table.columns[i].width = col_width

    # Style header row
    for i, header_text in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header_text

        # Header cell fill
        cell_fill = cell.fill
        cell_fill.solid()
        cell_fill.fore_color.rgb = RGBColor.from_string(colors['accent'])

        # Header text formatting
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor.from_string(colors['text'])
            paragraph.font.name = fonts['heading']
            paragraph.alignment = PP_ALIGN.CENTER

    # Style data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            if col_idx >= num_cols:
                break
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)

            # Alternating row colors
            cell_fill = cell.fill
            cell_fill.solid()
            if row_idx % 2 == 0:
                cell_fill.fore_color.rgb = RGBColor.from_string(colors['secondary'])
            else:
                cell_fill.fore_color.rgb = RGBColor.from_string(colors['bg'])

            # Data text formatting
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.color.rgb = RGBColor.from_string(colors['text'])
                paragraph.font.name = fonts['body']
                paragraph.alignment = PP_ALIGN.CENTER

    return slide
