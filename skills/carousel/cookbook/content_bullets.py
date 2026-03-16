"""Layout: Content Bullets — Title + bullet point list."""
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def add_slide(prs, title="", body="", colors=None, fonts=None, **kwargs):
    """Add a slide with title and bullet points.

    Classic presentation layout — heading at top with a list of bullet
    points below. Bullets can be passed as a newline-separated string in
    body, or as a list via kwargs['items'].

    Args:
        prs: python-pptx Presentation object
        title: Slide heading
        body: Bullet points separated by newlines (alternative to items)
        colors: dict with hex colors
        fonts: dict with font names
        **kwargs:
            items (list[str]): List of bullet point strings

    Returns:
        The created slide object
    """
    if not colors:
        colors = {
            'primary': '1a1a2e',
            'secondary': '16213e',
            'accent': 'e94560',
            'text': 'ffffff',
            'bg': '1a1a2e',
        }
    if not fonts:
        fonts = {'heading': 'Arial Black', 'body': 'Arial'}

    items = kwargs.get('items', None)
    if not items and body:
        items = [line.strip() for line in body.split('\n') if line.strip()]
    if not items:
        items = []

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(colors['bg'])

    # Title
    title_left = Inches(1)
    title_top = Inches(0.8)
    title_width = slide_width - Inches(2)
    title_height = Inches(1)

    txBox = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor.from_string(colors['text'])
    p.font.name = fonts['heading']
    p.alignment = PP_ALIGN.LEFT

    # Accent underline
    line = slide.shapes.add_shape(
        1, title_left, int(title_top + title_height), Inches(1.5), Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor.from_string(colors['accent'])
    line.line.fill.background()

    # Bullet points
    bullets_left = Inches(1.2)
    bullets_top = int(title_top + title_height + Inches(0.4))
    bullets_width = slide_width - Inches(2.4)
    bullets_height = slide_height - bullets_top - Inches(0.8)

    txBox2 = slide.shapes.add_textbox(bullets_left, int(bullets_top), bullets_width, int(bullets_height))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p2 = tf2.paragraphs[0]
        else:
            p2 = tf2.add_paragraph()

        p2.text = f"\u2022  {item}"
        p2.font.size = Pt(20)
        p2.font.color.rgb = RGBColor.from_string(colors['text'])
        p2.font.name = fonts['body']
        p2.alignment = PP_ALIGN.LEFT
        p2.space_after = Pt(12)

    return slide
