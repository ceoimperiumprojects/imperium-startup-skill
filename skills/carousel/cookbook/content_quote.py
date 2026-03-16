"""Layout: Content Quote — Large quote with attribution."""
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def add_slide(prs, title="", body="", colors=None, fonts=None, **kwargs):
    """Add a quote slide with large quotation marks and attribution.

    Displays a quote in large italic text with oversized quotation marks
    as visual decoration. Attribution line appears below in smaller text.

    Args:
        prs: python-pptx Presentation object
        title: Attribution (e.g., "— Steve Jobs")
        body: The quote text
        colors: dict with hex colors
        fonts: dict with font names
        **kwargs:
            attribution (str): Alternative to using title for attribution

    Returns:
        The created slide object
    """
    if not colors:
        colors = {
            'primary': '1a1a2e',
            'secondary': '16213e',
            'accent': 'e94560',
            'text': 'ffffff',
            'bg': '1a1a2e',
        }
    if not fonts:
        fonts = {'heading': 'Arial Black', 'body': 'Arial'}

    attribution = kwargs.get('attribution', title)

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(colors['bg'])

    # Large decorative opening quote mark
    quote_mark = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(2), Inches(1.5))
    qtf = quote_mark.text_frame
    qp = qtf.paragraphs[0]
    qp.text = "\u201C"
    qp.font.size = Pt(120)
    qp.font.color.rgb = RGBColor.from_string(colors['accent'])
    qp.font.name = fonts['heading']
    qp.alignment = PP_ALIGN.LEFT

    # Quote text
    quote_left = Inches(1.5)
    quote_top = int(slide_height * 0.28)
    quote_width = slide_width - Inches(3)
    quote_height = int(slide_height * 0.40)

    txBox = slide.shapes.add_textbox(quote_left, quote_top, quote_width, quote_height)
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = body
    p.font.size = Pt(28)
    p.font.italic = True
    p.font.color.rgb = RGBColor.from_string(colors['text'])
    p.font.name = fonts['body']
    p.alignment = PP_ALIGN.LEFT
    p.line_spacing = Pt(40)

    # Attribution
    if attribution:
        attr_top = int(quote_top + quote_height + Inches(0.3))
        attr_box = slide.shapes.add_textbox(
            quote_left, attr_top, quote_width, Inches(0.6)
        )
        atf = attr_box.text_frame

        ap = atf.paragraphs[0]
        # Add em-dash if not already present
        attr_text = attribution if attribution.startswith('\u2014') or attribution.startswith('--') else f"\u2014 {attribution}"
        ap.text = attr_text
        ap.font.size = Pt(18)
        ap.font.bold = True
        ap.font.color.rgb = RGBColor.from_string(colors['accent'])
        ap.font.name = fonts['body']
        ap.alignment = PP_ALIGN.LEFT

    return slide
