"""Layout: Image Full — Full-bleed image with text overlay."""
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
import os


def add_slide(prs, title="", body="", colors=None, fonts=None, **kwargs):
    """Add a full-bleed image slide with text overlay at the bottom.

    Image fills the entire slide. A semi-transparent dark strip at the
    bottom holds the title and body text. If no image is provided, uses
    a solid dark background.

    Args:
        prs: python-pptx Presentation object
        title: Overlay title text
        body: Overlay body text
        colors: dict with hex colors
        fonts: dict with font names
        **kwargs:
            image_path (str): Path to image file

    Returns:
        The created slide object
    """
    if not colors:
        colors = {
            'primary': '1a1a2e',
            'secondary': '16213e',
            'accent': 'e94560',
            'text': 'ffffff',
            'bg': '1a1a2e',
        }
    if not fonts:
        fonts = {'heading': 'Arial Black', 'body': 'Arial'}

    image_path = kwargs.get('image_path', None)

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Full-bleed image or solid background
    if image_path and os.path.exists(image_path):
        slide.shapes.add_picture(image_path, 0, 0, slide_width, slide_height)
    else:
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(colors['secondary'])

    # Bottom overlay strip (semi-transparent)
    strip_height = int(slide_height * 0.35)
    strip_top = slide_height - strip_height

    overlay = slide.shapes.add_shape(
        1, 0, int(strip_top), slide_width, strip_height
    )
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(0, 0, 0)
    overlay.line.fill.background()

    # Set transparency to 30%
    solidFill = overlay.fill._fill
    srgbClr = solidFill.find(qn('a:solidFill'))
    if srgbClr is not None:
        clr_elem = srgbClr.find(qn('a:srgbClr'))
        if clr_elem is not None:
            alpha = clr_elem.makeelement(qn('a:alpha'), {'val': '30000'})
            clr_elem.append(alpha)

    # Title on overlay
    if title:
        txBox = slide.shapes.add_textbox(
            Inches(0.8), int(strip_top + Inches(0.3)),
            slide_width - Inches(1.6), Inches(0.8)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(30)
        p.font.bold = True
        p.font.color.rgb = RGBColor.from_string(colors['text'])
        p.font.name = fonts['heading']
        p.alignment = PP_ALIGN.LEFT

    # Body on overlay
    if body:
        body_top = int(strip_top + Inches(1.2)) if title else int(strip_top + Inches(0.3))
        txBox2 = slide.shapes.add_textbox(
            Inches(0.8), body_top,
            slide_width - Inches(1.6), Inches(1)
        )
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = body
        p2.font.size = Pt(16)
        p2.font.color.rgb = RGBColor.from_string(colors['text'])
        p2.font.name = fonts['body']
        p2.alignment = PP_ALIGN.LEFT

    return slide
