"""Layout: Title Bold — Large bold title, left-aligned, with accent bar."""
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor


def add_slide(prs, title="", body="", colors=None, fonts=None, **kwargs):
    """Add a bold title slide with left-aligned text and vertical accent bar.

    High-impact title slide with oversized typography and a colored
    accent bar on the left edge. Great for making a strong first impression.

    Args:
        prs: python-pptx Presentation object
        title: Main title text (keep it short — 3-6 words)
        body: Supporting subtitle or tagline
        colors: dict with 'primary', 'secondary', 'accent', 'text', 'bg' hex colors
        fonts: dict with 'heading', 'body' font names
        **kwargs: Additional parameters (unused)

    Returns:
        The created slide object
    """
    if not colors:
        colors = {
            'primary': '1a1a2e',
            'secondary': '16213e',
            'accent': 'e94560',
            'text': 'ffffff',
            'bg': '1a1a2e',
        }
    if not fonts:
        fonts = {'heading': 'Arial Black', 'body': 'Arial'}

    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(colors['bg'])

    # Vertical accent bar on the left
    bar_width = Inches(0.15)
    bar = slide.shapes.add_shape(
        1,  # RECTANGLE
        Inches(0.8), Inches(1.5), int(bar_width), int(slide_height - Inches(3))
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor.from_string(colors['accent'])
    bar.line.fill.background()

    # Title — large, bold, left-aligned
    title_left = Inches(1.3)
    title_top = Inches(1.5)
    title_width = slide_width - Inches(2.5)
    title_height = Inches(2.5)

    txBox = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor.from_string(colors['text'])
    p.font.name = fonts['heading']
    p.alignment = PP_ALIGN.LEFT

    # Subtitle
    if body:
        sub_left = Inches(1.3)
        sub_top = Inches(4.2)
        sub_width = slide_width - Inches(2.5)
        sub_height = Inches(1.2)

        txBox2 = slide.shapes.add_textbox(sub_left, sub_top, sub_width, sub_height)
        tf2 = txBox2.text_frame
        tf2.word_wrap = True

        p2 = tf2.paragraphs[0]
        p2.text = body
        p2.font.size = Pt(20)
        p2.font.color.rgb = RGBColor.from_string(colors['text'])
        p2.font.name = fonts['body']
        p2.alignment = PP_ALIGN.LEFT

    return slide
