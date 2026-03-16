"""Layout: Title Image — Title over background image with dark overlay."""
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os


def add_slide(prs, title="", body="", colors=None, fonts=None, **kwargs):
    """Add a title slide with background image and semi-transparent overlay.

    Places an image as a full-bleed background with a dark overlay shape
    on top, then renders title and subtitle text over it. Falls back to
    a solid background if no image_path is provided.

    Args:
        prs: python-pptx Presentation object
        title: Main title text
        body: Subtitle text
        colors: dict with hex colors
        fonts: dict with font names
        **kwargs:
            image_path (str): Path to background image file

    Returns:
        The created slide object
    """
    if not colors:
        colors = {
            'primary': '1a1a2e',
            'secondary': '16213e',
            'accent': 'e94560',
            'text': 'ffffff',
            'bg': '1a1a2e',
        }
    if not fonts:
        fonts = {'heading': 'Arial Black', 'body': 'Arial'}

    image_path = kwargs.get('image_path', None)

    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Background image (if provided)
    if image_path and os.path.exists(image_path):
        slide.shapes.add_picture(image_path, 0, 0, slide_width, slide_height)
    else:
        # Solid background fallback
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(colors['secondary'])

    # Semi-transparent dark overlay (simulated with a dark shape)
    overlay = slide.shapes.add_shape(
        1,  # RECTANGLE
        0, 0, slide_width, slide_height
    )
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(0, 0, 0)
    # python-pptx does not natively support shape transparency via fill,
    # but we can set it via the XML
    from pptx.oxml.ns import qn
    solidFill = overlay.fill._fill
    srgbClr = solidFill.find(qn('a:solidFill')).find(qn('a:srgbClr'))
    if srgbClr is not None:
        alpha = srgbClr.makeelement(qn('a:alpha'), {'val': '50000'})
        srgbClr.append(alpha)
    overlay.line.fill.background()

    # Title
    title_left = Inches(1)
    title_top = int(slide_height * 0.35)
    title_width = slide_width - Inches(2)
    title_height = Inches(1.8)

    txBox = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor.from_string(colors['text'])
    p.font.name = fonts['heading']
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    if body:
        sub_top = int(title_top + title_height + Inches(0.2))
        sub_left = Inches(1.5)
        sub_width = slide_width - Inches(3)
        sub_height = Inches(1)

        txBox2 = slide.shapes.add_textbox(sub_left, sub_top, sub_width, sub_height)
        tf2 = txBox2.text_frame
        tf2.word_wrap = True

        p2 = tf2.paragraphs[0]
        p2.text = body
        p2.font.size = Pt(22)
        p2.font.color.rgb = RGBColor.from_string(colors['text'])
        p2.font.name = fonts['body']
        p2.alignment = PP_ALIGN.CENTER

    return slide
