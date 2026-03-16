"""Layout: Title Centered — Centered title + subtitle, full background color."""
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor


def add_slide(prs, title="", body="", colors=None, fonts=None, **kwargs):
    """Add a centered title slide to the presentation.

    Clean, centered layout with title and subtitle stacked vertically
    on a solid background. Ideal for cover slides and section openers.

    Args:
        prs: python-pptx Presentation object
        title: Main title text
        body: Subtitle text
        colors: dict with 'primary', 'secondary', 'accent', 'text', 'bg' hex colors
        fonts: dict with 'heading', 'body' font names
        **kwargs: Additional parameters (unused)

    Returns:
        The created slide object
    """
    if not colors:
        colors = {
            'primary': '1a1a2e',
            'secondary': '16213e',
            'accent': 'e94560',
            'text': 'ffffff',
            'bg': '1a1a2e',
        }
    if not fonts:
        fonts = {'heading': 'Arial Black', 'body': 'Arial'}

    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Full-slide background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(colors['bg'])

    # Title text box — centered horizontally, upper-center vertically
    title_left = Inches(1)
    title_top = slide_height * 0.30
    title_width = slide_width - Inches(2)
    title_height = Inches(1.5)

    txBox = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor.from_string(colors['text'])
    p.font.name = fonts['heading']
    p.alignment = PP_ALIGN.CENTER

    # Accent line below title
    line_width = Inches(2)
    line_left = int((slide_width - line_width) / 2)
    line_top = int(title_top + title_height + Inches(0.2))
    line = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        line_left, line_top, int(line_width), Pt(4)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor.from_string(colors['accent'])
    line.line.fill.background()

    # Subtitle text box
    if body:
        sub_left = Inches(1.5)
        sub_top = int(line_top + Inches(0.5))
        sub_width = slide_width - Inches(3)
        sub_height = Inches(1.2)

        txBox2 = slide.shapes.add_textbox(sub_left, sub_top, sub_width, sub_height)
        tf2 = txBox2.text_frame
        tf2.word_wrap = True

        p2 = tf2.paragraphs[0]
        p2.text = body
        p2.font.size = Pt(20)
        p2.font.color.rgb = RGBColor.from_string(colors['text'])
        p2.font.name = fonts['body']
        p2.alignment = PP_ALIGN.CENTER

    return slide
