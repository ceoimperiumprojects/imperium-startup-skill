"""Layout: Content Checklist — Checklist format with checkmarks."""
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def add_slide(prs, title="", body="", colors=None, fonts=None, **kwargs):
    """Add a checklist slide with checkmark icons next to each item.

    Displays a title and a vertical list of items, each prefixed with a
    colored checkmark. Supports marking items as checked or unchecked.

    Args:
        prs: python-pptx Presentation object
        title: Slide heading
        body: Checklist items separated by newlines
        colors: dict with hex colors
        fonts: dict with font names
        **kwargs:
            items (list[str]): List of checklist item strings
            checked (list[bool]): Which items are checked (default: all True)

    Returns:
        The created slide object
    """
    if not colors:
        colors = {
            'primary': '1a1a2e',
            'secondary': '16213e',
            'accent': 'e94560',
            'text': 'ffffff',
            'bg': '1a1a2e',
        }
    if not fonts:
        fonts = {'heading': 'Arial Black', 'body': 'Arial'}

    items = kwargs.get('items', None)
    checked = kwargs.get('checked', None)

    if not items and body:
        items = [line.strip() for line in body.split('\n') if line.strip()]
    if not items:
        items = []
    if not checked:
        checked = [True] * len(items)

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(colors['bg'])

    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(0.6), slide_width - Inches(2), Inches(0.9))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor.from_string(colors['text'])
    p.font.name = fonts['heading']
    p.alignment = PP_ALIGN.LEFT

    # Checklist items
    list_left = Inches(1.2)
    list_top = Inches(2.0)
    item_height = Inches(0.55)
    max_items = min(len(items), 10)  # Cap at 10 to fit on slide

    for i in range(max_items):
        item = items[i]
        is_checked = checked[i] if i < len(checked) else True
        y = int(list_top + i * item_height)

        # Checkmark or empty square
        check_char = "\u2705" if is_checked else "\u2B1C"
        check_box = slide.shapes.add_textbox(
            int(list_left), y, Inches(0.5), int(item_height)
        )
        ctf = check_box.text_frame
        cp = ctf.paragraphs[0]
        cp.text = check_char
        cp.font.size = Pt(18)
        cp.alignment = PP_ALIGN.CENTER

        # Item text
        item_box = slide.shapes.add_textbox(
            int(list_left + Inches(0.6)), y,
            slide_width - int(list_left) - Inches(1.8), int(item_height)
        )
        itf = item_box.text_frame
        itf.word_wrap = True
        ip = itf.paragraphs[0]
        ip.text = item
        ip.font.size = Pt(18)
        ip.font.color.rgb = RGBColor.from_string(colors['text'])
        ip.font.name = fonts['body']
        ip.alignment = PP_ALIGN.LEFT

    return slide
