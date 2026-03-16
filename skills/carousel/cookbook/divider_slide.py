"""Layout: Divider Slide — Section divider with large text."""
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def add_slide(prs, title="", body="", colors=None, fonts=None, **kwargs):
    """Add a section divider slide with large centered text.

    Full-color background with large section title and optional
    section number. Used to break a deck into logical sections and
    give the audience a visual pause.

    Args:
        prs: python-pptx Presentation object
        title: Section title (large, centered)
        body: Optional subtitle or section description
        colors: dict with hex colors
        fonts: dict with font names
        **kwargs:
            section_number (str): Section number to display (e.g., "01", "Part 2")

    Returns:
        The created slide object
    """
    if not colors:
        colors = {
            'primary': '1a1a2e',
            'secondary': '16213e',
            'accent': 'e94560',
            'text': 'ffffff',
            'bg': '1a1a2e',
        }
    if not fonts:
        fonts = {'heading': 'Arial Black', 'body': 'Arial'}

    section_number = kwargs.get('section_number', '')

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Accent-colored background (use accent, not bg, for visual break)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(colors['accent'])

    # Section number (large, faded)
    if section_number:
        num_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(2), Inches(1)
        )
        ntf = num_box.text_frame
        np = ntf.paragraphs[0]
        np.text = section_number
        np.font.size = Pt(60)
        np.font.bold = True
        np.font.color.rgb = RGBColor.from_string(colors['text'])
        np.font.name = fonts['heading']
        np.alignment = PP_ALIGN.LEFT

    # Section title — large and centered
    title_top = int(slide_height * 0.32)
    txBox = slide.shapes.add_textbox(
        Inches(1), title_top, slide_width - Inches(2), Inches(2)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = RGBColor.from_string(colors['text'])
    p.font.name = fonts['heading']
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    if body:
        sub_top = int(title_top + Inches(2.2))
        txBox2 = slide.shapes.add_textbox(
            Inches(1.5), sub_top, slide_width - Inches(3), Inches(0.8)
        )
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = body
        p2.font.size = Pt(20)
        p2.font.color.rgb = RGBColor.from_string(colors['text'])
        p2.font.name = fonts['body']
        p2.alignment = PP_ALIGN.CENTER

    return slide
