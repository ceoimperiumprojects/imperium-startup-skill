"""Layout: Content Timeline — Horizontal timeline with milestones."""
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def add_slide(prs, title="", body="", colors=None, fonts=None, **kwargs):
    """Add a horizontal timeline slide with milestone markers.

    Draws a horizontal line with circular milestone markers and labels.
    Each milestone has a date/label above and description below. Supports
    3-6 milestones for best visual results.

    Args:
        prs: python-pptx Presentation object
        title: Slide heading (e.g., "Product Roadmap")
        body: Unused — use milestones kwarg
        colors: dict with hex colors
        fonts: dict with font names
        **kwargs:
            milestones (list[dict]): Each dict has 'date' and 'label' keys
                Example: [{'date': 'Q1 2024', 'label': 'MVP Launch'}]

    Returns:
        The created slide object
    """
    if not colors:
        colors = {
            'primary': '1a1a2e',
            'secondary': '16213e',
            'accent': 'e94560',
            'text': 'ffffff',
            'bg': '1a1a2e',
        }
    if not fonts:
        fonts = {'heading': 'Arial Black', 'body': 'Arial'}

    milestones = kwargs.get('milestones', [])

    # Parse body as fallback: "date: label" per line
    if not milestones and body:
        for line in body.strip().split('\n'):
            line = line.strip()
            if ':' in line:
                parts = line.split(':', 1)
                milestones.append({'date': parts[0].strip(), 'label': parts[1].strip()})
            elif line:
                milestones.append({'date': '', 'label': line})

    if not milestones:
        milestones = [
            {'date': 'Start', 'label': 'Beginning'},
            {'date': 'Mid', 'label': 'Progress'},
            {'date': 'End', 'label': 'Complete'},
        ]

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(colors['bg'])

    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), slide_width - Inches(2), Inches(0.9))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor.from_string(colors['text'])
    p.font.name = fonts['heading']
    p.alignment = PP_ALIGN.LEFT

    # Timeline line
    num = len(milestones)
    line_left = Inches(1.2)
    line_right = slide_width - Inches(1.2)
    line_width = int(line_right - line_left)
    line_y = int(slide_height * 0.50)

    timeline = slide.shapes.add_shape(
        1, int(line_left), line_y, line_width, Pt(4)
    )
    timeline.fill.solid()
    timeline.fill.fore_color.rgb = RGBColor.from_string(colors['accent'])
    timeline.line.fill.background()

    # Milestone markers
    dot_size = Inches(0.3)
    spacing = line_width / max(num - 1, 1) if num > 1 else 0

    for i, ms in enumerate(milestones):
        if num == 1:
            cx = int(line_left + line_width / 2)
        else:
            cx = int(line_left + i * spacing)

        # Circle marker
        dot = slide.shapes.add_shape(
            9,  # OVAL
            int(cx - dot_size / 2), int(line_y - dot_size / 2 + Pt(2)),
            int(dot_size), int(dot_size)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = RGBColor.from_string(colors['accent'])
        dot.line.color.rgb = RGBColor.from_string(colors['text'])
        dot.line.width = Pt(2)

        # Date label above
        date_width = Inches(1.6)
        date_box = slide.shapes.add_textbox(
            int(cx - date_width / 2), int(line_y - Inches(0.9)),
            int(date_width), Inches(0.5)
        )
        dtf = date_box.text_frame
        dtf.word_wrap = True
        dp = dtf.paragraphs[0]
        dp.text = ms.get('date', '')
        dp.font.size = Pt(12)
        dp.font.bold = True
        dp.font.color.rgb = RGBColor.from_string(colors['accent'])
        dp.font.name = fonts['heading']
        dp.alignment = PP_ALIGN.CENTER

        # Label below
        label_width = Inches(1.8)
        label_box = slide.shapes.add_textbox(
            int(cx - label_width / 2), int(line_y + Inches(0.4)),
            int(label_width), Inches(0.8)
        )
        ltf = label_box.text_frame
        ltf.word_wrap = True
        lp = ltf.paragraphs[0]
        lp.text = ms.get('label', '')
        lp.font.size = Pt(13)
        lp.font.color.rgb = RGBColor.from_string(colors['text'])
        lp.font.name = fonts['body']
        lp.alignment = PP_ALIGN.CENTER

    return slide
