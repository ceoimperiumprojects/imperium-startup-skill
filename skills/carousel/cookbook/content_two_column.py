"""Layout: Content Two Column — Two-column text layout with title."""
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def add_slide(prs, title="", body="", colors=None, fonts=None, **kwargs):
    """Add a two-column text layout slide.

    Title spans the full width at top. Below it, content is split into
    two columns. Column content is passed via kwargs or by splitting
    body on '|||' separator.

    Args:
        prs: python-pptx Presentation object
        title: Slide heading
        body: Full body text — split on '|||' for left/right columns
        colors: dict with hex colors
        fonts: dict with font names
        **kwargs:
            left_title (str): Left column heading
            right_title (str): Right column heading
            left_body (str): Left column text
            right_body (str): Right column text

    Returns:
        The created slide object
    """
    if not colors:
        colors = {
            'primary': '1a1a2e',
            'secondary': '16213e',
            'accent': 'e94560',
            'text': 'ffffff',
            'bg': '1a1a2e',
        }
    if not fonts:
        fonts = {'heading': 'Arial Black', 'body': 'Arial'}

    left_title = kwargs.get('left_title', '')
    right_title = kwargs.get('right_title', '')
    left_body = kwargs.get('left_body', '')
    right_body = kwargs.get('right_body', '')

    # Fallback: split body on '|||'
    if not left_body and not right_body and body:
        parts = body.split('|||')
        left_body = parts[0].strip() if len(parts) > 0 else ''
        right_body = parts[1].strip() if len(parts) > 1 else ''

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(colors['bg'])

    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(0.6), slide_width - Inches(2), Inches(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor.from_string(colors['text'])
    p.font.name = fonts['heading']
    p.alignment = PP_ALIGN.LEFT

    # Calculate column dimensions
    margin = Inches(1)
    gutter = Inches(0.5)
    col_width = int((slide_width - 2 * margin - gutter) / 2)
    col_top = Inches(2.0)
    col_height = slide_height - col_top - Inches(0.8)

    # Left column
    left_box = slide.shapes.add_textbox(int(margin), int(col_top), col_width, int(col_height))
    ltf = left_box.text_frame
    ltf.word_wrap = True

    if left_title:
        lp = ltf.paragraphs[0]
        lp.text = left_title
        lp.font.size = Pt(20)
        lp.font.bold = True
        lp.font.color.rgb = RGBColor.from_string(colors['accent'])
        lp.font.name = fonts['heading']
        lp.space_after = Pt(8)

        lp2 = ltf.add_paragraph()
        lp2.text = left_body
        lp2.font.size = Pt(16)
        lp2.font.color.rgb = RGBColor.from_string(colors['text'])
        lp2.font.name = fonts['body']
    else:
        lp = ltf.paragraphs[0]
        lp.text = left_body
        lp.font.size = Pt(16)
        lp.font.color.rgb = RGBColor.from_string(colors['text'])
        lp.font.name = fonts['body']

    # Right column
    right_left = int(margin + col_width + gutter)
    right_box = slide.shapes.add_textbox(right_left, int(col_top), col_width, int(col_height))
    rtf = right_box.text_frame
    rtf.word_wrap = True

    if right_title:
        rp = rtf.paragraphs[0]
        rp.text = right_title
        rp.font.size = Pt(20)
        rp.font.bold = True
        rp.font.color.rgb = RGBColor.from_string(colors['accent'])
        rp.font.name = fonts['heading']
        rp.space_after = Pt(8)

        rp2 = rtf.add_paragraph()
        rp2.text = right_body
        rp2.font.size = Pt(16)
        rp2.font.color.rgb = RGBColor.from_string(colors['text'])
        rp2.font.name = fonts['body']
    else:
        rp = rtf.paragraphs[0]
        rp.text = right_body
        rp.font.size = Pt(16)
        rp.font.color.rgb = RGBColor.from_string(colors['text'])
        rp.font.name = fonts['body']

    # Vertical divider line between columns
    divider_x = int(margin + col_width + gutter / 2 - Pt(1))
    divider = slide.shapes.add_shape(
        1, divider_x, int(col_top), Pt(2), int(col_height)
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = RGBColor.from_string(colors['accent'])
    divider.line.fill.background()

    return slide
