"""Layout: Image Grid — 2x2 image grid with captions."""
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os


def add_slide(prs, title="", body="", colors=None, fonts=None, **kwargs):
    """Add a 2x2 image grid slide with optional captions.

    Four images arranged in a grid with a title at the top. Each cell
    can have a caption below it. If images are not provided, colored
    placeholders are used.

    Args:
        prs: python-pptx Presentation object
        title: Slide heading
        body: Unused
        colors: dict with hex colors
        fonts: dict with font names
        **kwargs:
            images (list[str]): List of 4 image paths
            captions (list[str]): List of 4 caption strings

    Returns:
        The created slide object
    """
    if not colors:
        colors = {
            'primary': '1a1a2e',
            'secondary': '16213e',
            'accent': 'e94560',
            'text': 'ffffff',
            'bg': '1a1a2e',
        }
    if not fonts:
        fonts = {'heading': 'Arial Black', 'body': 'Arial'}

    images = kwargs.get('images', [None, None, None, None])
    captions = kwargs.get('captions', ['', '', '', ''])

    # Pad to 4 items
    while len(images) < 4:
        images.append(None)
    while len(captions) < 4:
        captions.append('')

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(colors['bg'])

    # Title
    if title:
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), slide_width - Inches(1.6), Inches(0.7))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor.from_string(colors['text'])
        p.font.name = fonts['heading']
        p.alignment = PP_ALIGN.LEFT

    # Grid layout
    margin = Inches(0.8)
    gap = Inches(0.3)
    grid_top = Inches(1.4)
    cell_width = int((slide_width - 2 * margin - gap) / 2)
    cell_img_height = int((slide_height - grid_top - Inches(0.8) - gap - Inches(0.8)) / 2)
    caption_height = Inches(0.4)

    positions = [
        (int(margin), int(grid_top)),                                              # top-left
        (int(margin + cell_width + gap), int(grid_top)),                           # top-right
        (int(margin), int(grid_top + cell_img_height + caption_height + gap)),     # bottom-left
        (int(margin + cell_width + gap), int(grid_top + cell_img_height + caption_height + gap)),  # bottom-right
    ]

    for i, (x, y) in enumerate(positions):
        img_path = images[i] if i < len(images) else None
        caption = captions[i] if i < len(captions) else ''

        if img_path and os.path.exists(img_path):
            slide.shapes.add_picture(img_path, x, y, cell_width, cell_img_height)
        else:
            # Placeholder rectangle
            ph = slide.shapes.add_shape(1, x, y, cell_width, cell_img_height)
            ph.fill.solid()
            ph.fill.fore_color.rgb = RGBColor.from_string(colors['secondary'])
            ph.line.color.rgb = RGBColor.from_string(colors['accent'])
            ph.line.width = Pt(1)

            # Placeholder text
            ph_txt = slide.shapes.add_textbox(
                int(x + cell_width * 0.3), int(y + cell_img_height * 0.4),
                int(cell_width * 0.4), Inches(0.4)
            )
            ptf = ph_txt.text_frame
            pp = ptf.paragraphs[0]
            pp.text = f"[IMG {i+1}]"
            pp.font.size = Pt(12)
            pp.font.color.rgb = RGBColor(100, 100, 100)
            pp.font.name = fonts['body']
            pp.alignment = PP_ALIGN.CENTER

        # Caption below image
        if caption:
            cap_box = slide.shapes.add_textbox(
                x, int(y + cell_img_height + Pt(4)), cell_width, int(caption_height)
            )
            ctf = cap_box.text_frame
            ctf.word_wrap = True
            cp = ctf.paragraphs[0]
            cp.text = caption
            cp.font.size = Pt(11)
            cp.font.color.rgb = RGBColor.from_string(colors['text'])
            cp.font.name = fonts['body']
            cp.alignment = PP_ALIGN.CENTER

    return slide
