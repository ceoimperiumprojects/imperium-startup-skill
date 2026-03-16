"""Layout: Title Gradient — Title with two-tone gradient background."""
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn


def add_slide(prs, title="", body="", colors=None, fonts=None, **kwargs):
    """Add a title slide with a two-color gradient background.

    Creates a diagonal gradient from primary to secondary color,
    with centered title and subtitle. Elegant opener for any deck.

    Args:
        prs: python-pptx Presentation object
        title: Main title text
        body: Subtitle text
        colors: dict with hex colors
        fonts: dict with font names
        **kwargs: Additional parameters (unused)

    Returns:
        The created slide object
    """
    if not colors:
        colors = {
            'primary': '1a1a2e',
            'secondary': '16213e',
            'accent': 'e94560',
            'text': 'ffffff',
            'bg': '1a1a2e',
        }
    if not fonts:
        fonts = {'heading': 'Arial Black', 'body': 'Arial'}

    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Create gradient background via a full-slide shape
    grad_shape = slide.shapes.add_shape(
        1,  # RECTANGLE
        0, 0, slide_width, slide_height
    )
    grad_shape.line.fill.background()

    # Set gradient fill via XML
    spPr = grad_shape._element.spPr
    # Remove any existing fill
    for child in list(spPr):
        if child.tag.endswith('}solidFill') or child.tag.endswith('}gradFill'):
            spPr.remove(child)

    gradFill = spPr.makeelement(qn('a:gradFill'), {'rotWithShape': '1'})
    gsLst = gradFill.makeelement(qn('a:gsLst'), {})

    # Stop 1 — primary color at 0%
    gs1 = gsLst.makeelement(qn('a:gs'), {'pos': '0'})
    srgb1 = gs1.makeelement(qn('a:srgbClr'), {'val': colors['primary']})
    gs1.append(srgb1)
    gsLst.append(gs1)

    # Stop 2 — secondary color at 100%
    gs2 = gsLst.makeelement(qn('a:gs'), {'pos': '100000'})
    srgb2 = gs2.makeelement(qn('a:srgbClr'), {'val': colors['secondary']})
    gs2.append(srgb2)
    gsLst.append(gs2)

    gradFill.append(gsLst)

    # Linear gradient — diagonal (45 degrees)
    lin = gradFill.makeelement(qn('a:lin'), {'ang': '2700000', 'scaled': '1'})
    gradFill.append(lin)

    spPr.append(gradFill)

    # Title
    title_left = Inches(1)
    title_top = int(slide_height * 0.30)
    title_width = slide_width - Inches(2)
    title_height = Inches(1.8)

    txBox = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor.from_string(colors['text'])
    p.font.name = fonts['heading']
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    if body:
        sub_top = int(title_top + title_height + Inches(0.3))
        txBox2 = slide.shapes.add_textbox(
            Inches(1.5), sub_top, slide_width - Inches(3), Inches(1)
        )
        tf2 = txBox2.text_frame
        tf2.word_wrap = True

        p2 = tf2.paragraphs[0]
        p2.text = body
        p2.font.size = Pt(22)
        p2.font.color.rgb = RGBColor.from_string(colors['text'])
        p2.font.name = fonts['body']
        p2.alignment = PP_ALIGN.CENTER

    return slide
