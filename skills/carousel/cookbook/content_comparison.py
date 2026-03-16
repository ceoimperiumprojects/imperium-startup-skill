"""Layout: Content Comparison — Side-by-side comparison (vs layout)."""
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def add_slide(prs, title="", body="", colors=None, fonts=None, **kwargs):
    """Add a comparison slide with two sides (before/after, us/them, etc.).

    Two-panel layout with a 'VS' or divider in the center. Each side
    has its own heading and content. Useful for before/after, competitor
    comparisons, or contrasting approaches.

    Args:
        prs: python-pptx Presentation object
        title: Slide heading (e.g., "Traditional vs Modern")
        body: Unused — use kwargs instead
        colors: dict with hex colors
        fonts: dict with font names
        **kwargs:
            left_title (str): Left side heading
            right_title (str): Right side heading
            left_items (list[str]): Left side bullet points
            right_items (list[str]): Right side bullet points
            left_color (str): Override color for left panel bg
            right_color (str): Override color for right panel bg

    Returns:
        The created slide object
    """
    if not colors:
        colors = {
            'primary': '1a1a2e',
            'secondary': '16213e',
            'accent': 'e94560',
            'text': 'ffffff',
            'bg': '1a1a2e',
        }
    if not fonts:
        fonts = {'heading': 'Arial Black', 'body': 'Arial'}

    left_title = kwargs.get('left_title', 'Before')
    right_title = kwargs.get('right_title', 'After')
    left_items = kwargs.get('left_items', [])
    right_items = kwargs.get('right_items', [])

    # Parse body as fallback: "left content ||| right content"
    if not left_items and not right_items and body:
        parts = body.split('|||')
        if len(parts) >= 2:
            left_items = [l.strip() for l in parts[0].strip().split('\n') if l.strip()]
            right_items = [r.strip() for r in parts[1].strip().split('\n') if r.strip()]

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(colors['bg'])

    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), slide_width - Inches(2), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = RGBColor.from_string(colors['text'])
    p.font.name = fonts['heading']
    p.alignment = PP_ALIGN.CENTER

    # Panel dimensions
    margin = Inches(0.8)
    panel_gap = Inches(0.8)
    panel_width = int((slide_width - 2 * margin - panel_gap) / 2)
    panel_top = Inches(1.8)
    panel_height = slide_height - panel_top - Inches(0.6)

    # Left panel background
    left_panel = slide.shapes.add_shape(
        1, int(margin), int(panel_top), panel_width, int(panel_height)
    )
    left_color = kwargs.get('left_color', colors['secondary'])
    left_panel.fill.solid()
    left_panel.fill.fore_color.rgb = RGBColor.from_string(left_color)
    left_panel.line.fill.background()

    # Left title
    lt_box = slide.shapes.add_textbox(
        int(margin + Inches(0.3)), int(panel_top + Inches(0.3)),
        int(panel_width - Inches(0.6)), Inches(0.6)
    )
    ltf = lt_box.text_frame
    ltp = ltf.paragraphs[0]
    ltp.text = left_title
    ltp.font.size = Pt(22)
    ltp.font.bold = True
    ltp.font.color.rgb = RGBColor.from_string(colors['accent'])
    ltp.font.name = fonts['heading']
    ltp.alignment = PP_ALIGN.CENTER

    # Left items
    li_box = slide.shapes.add_textbox(
        int(margin + Inches(0.4)), int(panel_top + Inches(1.1)),
        int(panel_width - Inches(0.8)), int(panel_height - Inches(1.5))
    )
    litf = li_box.text_frame
    litf.word_wrap = True
    for i, item in enumerate(left_items):
        lp = litf.paragraphs[0] if i == 0 else litf.add_paragraph()
        lp.text = f"\u2022  {item}"
        lp.font.size = Pt(16)
        lp.font.color.rgb = RGBColor.from_string(colors['text'])
        lp.font.name = fonts['body']
        lp.space_after = Pt(8)

    # VS circle in the center
    vs_size = Inches(0.8)
    vs_x = int((slide_width - vs_size) / 2)
    vs_y = int(panel_top + panel_height / 2 - vs_size / 2)
    vs_circle = slide.shapes.add_shape(
        9,  # MSO_SHAPE.OVAL
        vs_x, vs_y, int(vs_size), int(vs_size)
    )
    vs_circle.fill.solid()
    vs_circle.fill.fore_color.rgb = RGBColor.from_string(colors['accent'])
    vs_circle.line.fill.background()

    vs_tf = vs_circle.text_frame
    vs_tf.word_wrap = False
    vsp = vs_tf.paragraphs[0]
    vsp.text = "VS"
    vsp.font.size = Pt(14)
    vsp.font.bold = True
    vsp.font.color.rgb = RGBColor.from_string(colors['text'])
    vsp.font.name = fonts['heading']
    vsp.alignment = PP_ALIGN.CENTER

    # Right panel background
    right_x = int(margin + panel_width + panel_gap)
    right_panel = slide.shapes.add_shape(
        1, right_x, int(panel_top), panel_width, int(panel_height)
    )
    right_color_val = kwargs.get('right_color', colors['secondary'])
    right_panel.fill.solid()
    right_panel.fill.fore_color.rgb = RGBColor.from_string(right_color_val)
    right_panel.line.fill.background()

    # Right title
    rt_box = slide.shapes.add_textbox(
        int(right_x + Inches(0.3)), int(panel_top + Inches(0.3)),
        int(panel_width - Inches(0.6)), Inches(0.6)
    )
    rtf = rt_box.text_frame
    rtp = rtf.paragraphs[0]
    rtp.text = right_title
    rtp.font.size = Pt(22)
    rtp.font.bold = True
    rtp.font.color.rgb = RGBColor.from_string(colors['accent'])
    rtp.font.name = fonts['heading']
    rtp.alignment = PP_ALIGN.CENTER

    # Right items
    ri_box = slide.shapes.add_textbox(
        int(right_x + Inches(0.4)), int(panel_top + Inches(1.1)),
        int(panel_width - Inches(0.8)), int(panel_height - Inches(1.5))
    )
    ritf = ri_box.text_frame
    ritf.word_wrap = True
    for i, item in enumerate(right_items):
        rp = ritf.paragraphs[0] if i == 0 else ritf.add_paragraph()
        rp.text = f"\u2022  {item}"
        rp.font.size = Pt(16)
        rp.font.color.rgb = RGBColor.from_string(colors['text'])
        rp.font.name = fonts['body']
        rp.space_after = Pt(8)

    return slide
