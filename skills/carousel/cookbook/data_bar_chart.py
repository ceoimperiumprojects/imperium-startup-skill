"""Layout: Data Bar Chart — Simple bar chart using python-pptx charts."""
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData


def add_slide(prs, title="", body="", colors=None, fonts=None, **kwargs):
    """Add a bar chart slide using python-pptx native charting.

    Creates a column chart with category labels and data values.
    Supports single or multiple series.

    Args:
        prs: python-pptx Presentation object
        title: Chart title / slide heading
        body: Unused
        colors: dict with hex colors
        fonts: dict with font names
        **kwargs:
            categories (list[str]): Category labels (x-axis)
            values (list[float]): Data values for single series
            series_data (dict[str, list[float]]): Named series for multi-series
                Example: {'Revenue': [10, 20, 30], 'Costs': [5, 10, 15]}

    Returns:
        The created slide object
    """
    if not colors:
        colors = {
            'primary': '1a1a2e',
            'secondary': '16213e',
            'accent': 'e94560',
            'text': 'ffffff',
            'bg': '1a1a2e',
        }
    if not fonts:
        fonts = {'heading': 'Arial Black', 'body': 'Arial'}

    categories = kwargs.get('categories', ['Q1', 'Q2', 'Q3', 'Q4'])
    values = kwargs.get('values', None)
    series_data = kwargs.get('series_data', None)

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(colors['bg'])

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), slide_width - Inches(1.6), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor.from_string(colors['text'])
    p.font.name = fonts['heading']
    p.alignment = PP_ALIGN.LEFT

    # Build chart data
    chart_data = CategoryChartData()
    chart_data.categories = categories

    if series_data:
        for series_name, series_values in series_data.items():
            chart_data.add_series(series_name, series_values)
    elif values:
        chart_data.add_series('Data', values)
    else:
        chart_data.add_series('Data', [25, 45, 60, 80])

    # Add chart
    chart_left = Inches(0.8)
    chart_top = Inches(1.6)
    chart_width = slide_width - Inches(1.6)
    chart_height = slide_height - Inches(2.2)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        int(chart_left), int(chart_top),
        int(chart_width), int(chart_height),
        chart_data
    )

    chart = chart_frame.chart
    chart.has_legend = bool(series_data and len(series_data) > 1)
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False

    # Style the chart
    plot = chart.plots[0]
    plot.gap_width = 100

    # Color the bars with accent color
    for i, series in enumerate(plot.series):
        fill = series.format.fill
        fill.solid()
        if i == 0:
            fill.fore_color.rgb = RGBColor.from_string(colors['accent'])
        else:
            # Alternate colors for multi-series
            fill.fore_color.rgb = RGBColor.from_string(colors['primary'])

    return slide
