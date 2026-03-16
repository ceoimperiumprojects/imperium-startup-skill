"""Layout: Image Right — Text on the left, image on the right."""
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os


def add_slide(prs, title="", body="", colors=None, fonts=None, **kwargs):
    """Add a slide with text on the left and image on the right half.

    Mirror of image_left layout. Text content on the left 55%, image
    on the right 45%. Good for alternating with image_left for visual rhythm.

    Args:
        prs: python-pptx Presentation object
        title: Text heading (left side)
        body: Text body content (left side)
        colors: dict with hex colors
        fonts: dict with font names
        **kwargs:
            image_path (str): Path to image file

    Returns:
        The created slide object
    """
    if not colors:
        colors = {
            'primary': '1a1a2e',
            'secondary': '16213e',
            'accent': 'e94560',
            'text': 'ffffff',
            'bg': '1a1a2e',
        }
    if not fonts:
        fonts = {'heading': 'Arial Black', 'body': 'Arial'}

    image_path = kwargs.get('image_path', None)

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(colors['bg'])

    # Right side — image or placeholder
    img_width = int(slide_width * 0.45)
    img_left = slide_width - img_width

    if image_path and os.path.exists(image_path):
        slide.shapes.add_picture(image_path, int(img_left), 0, img_width, slide_height)
    else:
        placeholder = slide.shapes.add_shape(
            1, int(img_left), 0, img_width, slide_height
        )
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = RGBColor.from_string(colors['secondary'])
        placeholder.line.fill.background()

        ph_text = slide.shapes.add_textbox(
            int(img_left + img_width * 0.2), int(slide_height * 0.45),
            int(img_width * 0.6), Inches(0.5)
        )
        phtf = ph_text.text_frame
        php = phtf.paragraphs[0]
        php.text = "[IMAGE]"
        php.font.size = Pt(14)
        php.font.color.rgb = RGBColor(100, 100, 100)
        php.font.name = fonts['body']
        php.alignment = PP_ALIGN.CENTER

    # Left side — text content
    text_left = Inches(0.8)
    text_width = int(img_left - Inches(1.4))

    # Title
    txBox = slide.shapes.add_textbox(
        int(text_left), Inches(1.2), text_width, Inches(1.2)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = RGBColor.from_string(colors['text'])
    p.font.name = fonts['heading']
    p.alignment = PP_ALIGN.LEFT

    # Accent bar
    accent_bar = slide.shapes.add_shape(
        1, int(text_left), Inches(2.5), Inches(1.2), Pt(3)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = RGBColor.from_string(colors['accent'])
    accent_bar.line.fill.background()

    # Body
    if body:
        txBox2 = slide.shapes.add_textbox(
            int(text_left), Inches(2.9), text_width, slide_height - Inches(3.8)
        )
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = body
        p2.font.size = Pt(16)
        p2.font.color.rgb = RGBColor.from_string(colors['text'])
        p2.font.name = fonts['body']
        p2.alignment = PP_ALIGN.LEFT
        p2.line_spacing = Pt(24)

    return slide
