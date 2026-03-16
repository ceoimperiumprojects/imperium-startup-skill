"""Layout: Data Pie Chart — Pie chart with labels."""
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData


def add_slide(prs, title="", body="", colors=None, fonts=None, **kwargs):
    """Add a pie chart slide with data labels.

    Creates a pie chart showing proportional data. Labels show
    category name and percentage. Best with 3-6 categories.

    Args:
        prs: python-pptx Presentation object
        title: Chart title / slide heading
        body: Optional footnote or context text below chart
        colors: dict with hex colors
        fonts: dict with font names
        **kwargs:
            categories (list[str]): Slice labels
            values (list[float]): Slice values

    Returns:
        The created slide object
    """
    if not colors:
        colors = {
            'primary': '1a1a2e',
            'secondary': '16213e',
            'accent': 'e94560',
            'text': 'ffffff',
            'bg': '1a1a2e',
        }
    if not fonts:
        fonts = {'heading': 'Arial Black', 'body': 'Arial'}

    categories = kwargs.get('categories', ['Segment A', 'Segment B', 'Segment C'])
    values = kwargs.get('values', [45, 30, 25])

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(colors['bg'])

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), slide_width - Inches(1.6), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor.from_string(colors['text'])
    p.font.name = fonts['heading']
    p.alignment = PP_ALIGN.LEFT

    # Chart data
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('Data', values)

    # Add pie chart
    chart_left = Inches(1.5)
    chart_top = Inches(1.6)
    chart_width = slide_width - Inches(3)
    chart_height = slide_height - Inches(2.8)

    if body:
        chart_height = slide_height - Inches(3.6)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE,
        int(chart_left), int(chart_top),
        int(chart_width), int(chart_height),
        chart_data
    )

    chart = chart_frame.chart
    chart.has_legend = True
    chart.legend.position = 2  # RIGHT
    chart.legend.include_in_layout = False

    # Data labels — show percentage
    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.show_percentage = True
    data_labels.show_category_name = False
    data_labels.show_value = False
    data_labels.font.size = Pt(11)
    data_labels.font.color.rgb = RGBColor.from_string(colors['text'])

    # Color pie slices
    pie_colors = [colors['accent'], colors['primary'], colors['secondary'],
                  '4a90d9', '50c878', 'f5a623']
    series = plot.series[0]
    for i in range(len(categories)):
        point = series.points[i]
        point.format.fill.solid()
        color_hex = pie_colors[i % len(pie_colors)]
        point.format.fill.fore_color.rgb = RGBColor.from_string(color_hex)

    # Footnote / context
    if body:
        note_box = slide.shapes.add_textbox(
            Inches(0.8), int(slide_height - Inches(0.8)),
            slide_width - Inches(1.6), Inches(0.5)
        )
        ntf = note_box.text_frame
        ntf.word_wrap = True
        np = ntf.paragraphs[0]
        np.text = body
        np.font.size = Pt(12)
        np.font.color.rgb = RGBColor.from_string(colors['text'])
        np.font.name = fonts['body']
        np.alignment = PP_ALIGN.CENTER

    return slide
