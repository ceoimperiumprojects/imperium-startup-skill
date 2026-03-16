"""Layout: Content Stats — Big number/statistic with context."""
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def add_slide(prs, title="", body="", colors=None, fonts=None, **kwargs):
    """Add a statistics slide with a big number and supporting context.

    Displays a large statistic (number, percentage, metric) prominently
    in the center with a label above and context below. Great for
    traction slides and data highlights.

    Args:
        prs: python-pptx Presentation object
        title: Label above the stat (e.g., "Monthly Active Users")
        body: Context or explanation below the stat
        colors: dict with hex colors
        fonts: dict with font names
        **kwargs:
            stat (str): The big number/stat to display (e.g., "10M+", "340%", "$2.5M")
            stats (list[dict]): Multiple stats, each with 'value' and 'label' keys

    Returns:
        The created slide object
    """
    if not colors:
        colors = {
            'primary': '1a1a2e',
            'secondary': '16213e',
            'accent': 'e94560',
            'text': 'ffffff',
            'bg': '1a1a2e',
        }
    if not fonts:
        fonts = {'heading': 'Arial Black', 'body': 'Arial'}

    stat = kwargs.get('stat', '')
    stats = kwargs.get('stats', None)

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(colors['bg'])

    if stats and len(stats) > 1:
        # Multiple stats layout — evenly spaced across slide
        num_stats = len(stats)
        stat_width = int((slide_width - Inches(2)) / num_stats)

        for i, s in enumerate(stats):
            x = int(Inches(1) + i * stat_width)
            y_val = int(slide_height * 0.30)

            # Value
            val_box = slide.shapes.add_textbox(x, y_val, stat_width, Inches(1.2))
            vtf = val_box.text_frame
            vtf.word_wrap = True
            vp = vtf.paragraphs[0]
            vp.text = str(s.get('value', ''))
            vp.font.size = Pt(48)
            vp.font.bold = True
            vp.font.color.rgb = RGBColor.from_string(colors['accent'])
            vp.font.name = fonts['heading']
            vp.alignment = PP_ALIGN.CENTER

            # Label
            lbl_box = slide.shapes.add_textbox(x, int(y_val + Inches(1.3)), stat_width, Inches(0.6))
            ltf = lbl_box.text_frame
            ltf.word_wrap = True
            lp = ltf.paragraphs[0]
            lp.text = str(s.get('label', ''))
            lp.font.size = Pt(16)
            lp.font.color.rgb = RGBColor.from_string(colors['text'])
            lp.font.name = fonts['body']
            lp.alignment = PP_ALIGN.CENTER
    else:
        # Single stat layout
        # Label above
        if title:
            label_box = slide.shapes.add_textbox(
                Inches(1), int(slide_height * 0.20),
                slide_width - Inches(2), Inches(0.6)
            )
            ltf = label_box.text_frame
            lp = ltf.paragraphs[0]
            lp.text = title.upper()
            lp.font.size = Pt(16)
            lp.font.bold = True
            lp.font.color.rgb = RGBColor.from_string(colors['text'])
            lp.font.name = fonts['body']
            lp.alignment = PP_ALIGN.CENTER

        # Big number
        stat_text = stat if stat else body
        stat_box = slide.shapes.add_textbox(
            Inches(1), int(slide_height * 0.32),
            slide_width - Inches(2), Inches(2)
        )
        stf = stat_box.text_frame
        stf.word_wrap = True
        sp = stf.paragraphs[0]
        sp.text = stat_text
        sp.font.size = Pt(72)
        sp.font.bold = True
        sp.font.color.rgb = RGBColor.from_string(colors['accent'])
        sp.font.name = fonts['heading']
        sp.alignment = PP_ALIGN.CENTER

        # Context below
        if body and stat:
            ctx_box = slide.shapes.add_textbox(
                Inches(1.5), int(slide_height * 0.62),
                slide_width - Inches(3), Inches(1)
            )
            ctf = ctx_box.text_frame
            ctf.word_wrap = True
            cp = ctf.paragraphs[0]
            cp.text = body
            cp.font.size = Pt(18)
            cp.font.color.rgb = RGBColor.from_string(colors['text'])
            cp.font.name = fonts['body']
            cp.alignment = PP_ALIGN.CENTER

    return slide
